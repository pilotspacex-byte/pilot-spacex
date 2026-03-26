"""OfficeExtractionService — synchronous extraction of DOCX, XLSX, and PPTX bytes to markdown.

Converts Office document bytes into structured markdown text suitable for AI context
building and knowledge graph ingestion.

Feature: 020 — Office Document Extraction (Phase 41)
Requirements: OFFICE-01, OFFICE-02, OFFICE-03, OFFICE-04
"""

from __future__ import annotations

import io
from dataclasses import dataclass

from pilot_space.domain.exceptions import ValidationError

_MAX_ROWS_PER_SHEET = 500

_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

_SENTINEL_TEXT = "[Office file is password-protected and could not be extracted.]"


@dataclass(frozen=True)
class ExtractionResult:
    """Immutable result of an Office document extraction."""

    text: str
    metadata: dict[str, object]


class OfficeExtractionService:
    """Extracts text content from DOCX, XLSX, and PPTX files as structured markdown.

    This service is synchronous and stateless — safe to use as a singleton in the DI
    container (no session, no database, no async I/O).

    Usage::

        svc = OfficeExtractionService()
        result = svc.extract(data, mime_type, filename)
        print(result.text)
    """

    OFFICE_MIME_TYPES: frozenset[str] = frozenset(
        {
            _DOCX_MIME,
            _XLSX_MIME,
            _PPTX_MIME,
        }
    )

    def extract(self, data: bytes, mime_type: str, filename: str) -> ExtractionResult:
        """Extract text from Office document bytes.

        Args:
            data: Raw bytes of the Office document.
            mime_type: MIME type string used to route to the correct extractor.
            filename: Original filename (used for logging only, not for MIME detection).

        Returns:
            ExtractionResult with extracted markdown text and format-specific metadata.

        Raises:
            ValidationError: If mime_type is not one of the three supported Office MIME types.
        """
        if mime_type == _DOCX_MIME:
            return self._extract_docx(data)
        if mime_type == _XLSX_MIME:
            return self._extract_xlsx(data)
        if mime_type == _PPTX_MIME:
            return self._extract_pptx(data)
        msg = f"Unsupported MIME type for Office extraction: {mime_type}"
        raise ValidationError(msg)

    # ------------------------------------------------------------------
    # OFFICE-01: DOCX extraction via mammoth
    # ------------------------------------------------------------------

    def _extract_docx(self, data: bytes) -> ExtractionResult:
        """Convert DOCX bytes to markdown using mammoth."""
        try:
            import mammoth  # lazy import — not loaded on every module import

            def _image_handler(image: object) -> dict[str, str]:  # type: ignore[type-arg]
                """Replace embedded images with an [Image] placeholder."""
                return {"alt_text": "[Image]"}

            convert_image = mammoth.images.img_element(_image_handler)
            result = mammoth.convert_to_markdown(
                io.BytesIO(data),
                convert_image=convert_image,
            )
            text: str = result.value or ""
            word_count = len(text.split()) if text.strip() else 0
            return ExtractionResult(
                text=text,
                metadata={"word_count": word_count},
            )
        except Exception:
            return ExtractionResult(text=_SENTINEL_TEXT, metadata={})

    # ------------------------------------------------------------------
    # OFFICE-02: XLSX extraction via openpyxl
    # ------------------------------------------------------------------

    def _extract_xlsx(self, data: bytes) -> ExtractionResult:
        """Convert XLSX bytes to GFM markdown tables using openpyxl."""
        try:
            import openpyxl  # lazy import

            wb = openpyxl.load_workbook(
                io.BytesIO(data),
                read_only=True,
                data_only=True,
            )
            try:
                sections: list[str] = []
                sheet_count = len(wb.sheetnames)

                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    section_lines: list[str] = [f"## Sheet: {sheet_name}"]

                    rows_iter = ws.iter_rows(values_only=True)

                    # Skip leading blank rows; first non-blank row is the header.
                    header_row = None
                    for candidate in rows_iter:
                        if not all(v is None for v in candidate):
                            header_row = candidate
                            break

                    if header_row is None:
                        # Sheet is entirely blank — no table
                        sections.append("\n".join(section_lines))
                        continue

                    header_cells = [str(cell) if cell is not None else "" for cell in header_row]
                    section_lines.append("| " + " | ".join(header_cells) + " |")
                    section_lines.append("| " + " | ".join("---" for _ in header_cells) + " |")

                    data_row_count = 0
                    total_rows = 0
                    for row in rows_iter:
                        total_rows += 1
                        if data_row_count < _MAX_ROWS_PER_SHEET:
                            cells = [str(cell) if cell is not None else "" for cell in row]
                            section_lines.append("| " + " | ".join(cells) + " |")
                            data_row_count += 1

                    if total_rows > _MAX_ROWS_PER_SHEET:
                        section_lines.append(
                            f"*[Truncated: showing {_MAX_ROWS_PER_SHEET} of {total_rows} rows]*"
                        )

                    sections.append("\n".join(section_lines))

            finally:
                wb.close()

            text = "\n\n".join(sections)
            return ExtractionResult(
                text=text,
                metadata={"sheet_count": sheet_count},
            )
        except Exception:
            return ExtractionResult(text=_SENTINEL_TEXT, metadata={})

    # ------------------------------------------------------------------
    # OFFICE-03: PPTX extraction via python-pptx
    # ------------------------------------------------------------------

    def _extract_pptx(self, data: bytes) -> ExtractionResult:
        """Convert PPTX bytes to markdown with per-slide sections using python-pptx."""
        try:
            from pptx import Presentation  # lazy import

            prs = Presentation(io.BytesIO(data))
            sections: list[str] = []
            slide_count = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides, start=1):
                title_text: str | None = None

                # Find the title shape: prefer shape named "title*", fall back to
                # placeholder index 0, then first text shape.
                for shape in slide.shapes:
                    if shape.name.lower().startswith("title"):
                        if shape.has_text_frame:
                            title_text = shape.text_frame.text.strip()  # type: ignore[union-attr]
                            break
                    # Shape type 13 == FREEFORM which some decks use for title areas
                    # (secondary detection).

                if title_text is None:
                    # Try placeholder idx 0 (standard title placeholder)
                    try:
                        ph = slide.placeholders[0]
                        if ph.has_text_frame:
                            title_text = ph.text_frame.text.strip()  # type: ignore[union-attr]
                    except (KeyError, IndexError):
                        pass

                display_title = title_text if title_text else "(Untitled)"
                slide_lines: list[str] = [f"## Slide {slide_num}: {display_title}"]

                # Extract text from all non-title shapes
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    if shape.name.lower().startswith("title"):
                        continue  # already captured as title
                    for para in shape.text_frame.paragraphs:  # type: ignore[union-attr]
                        para_text = para.text.strip()
                        if para_text:
                            slide_lines.append(f"- {para_text}")

                # Speaker notes
                if slide.has_notes_slide:
                    notes_tf = slide.notes_slide.notes_text_frame
                    if notes_tf is not None:
                        notes_text = notes_tf.text.strip()
                        if notes_text:
                            slide_lines.append(f"> {notes_text}")

                sections.append("\n".join(slide_lines))

            text = "\n\n".join(sections)
            return ExtractionResult(
                text=text,
                metadata={"slide_count": slide_count},
            )
        except Exception:
            return ExtractionResult(text=_SENTINEL_TEXT, metadata={})
