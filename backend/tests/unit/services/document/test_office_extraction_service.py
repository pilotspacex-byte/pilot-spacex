"""Unit tests for OfficeExtractionService — OFFICE-01 through OFFICE-04.

Fixtures are generated in-memory using the same libraries the service uses.
No external fixture files required — pytest generates minimal valid Office
documents via mammoth/openpyxl/python-pptx programmatic APIs.

Feature: 020 — Office Document Extraction (Phase 41)
Requirements: OFFICE-01, OFFICE-02, OFFICE-03, OFFICE-04
"""

from __future__ import annotations

import io

import pytest

from pilot_space.application.services.document.office_extraction_service import (
    ExtractionResult,
    OfficeExtractionService,
)
from pilot_space.domain.exceptions import ValidationError

_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@pytest.fixture(scope="module")
def svc() -> OfficeExtractionService:
    return OfficeExtractionService()


def _build_minimal_docx(heading: str, body: str) -> bytes:
    """Build a minimal valid DOCX (OOXML ZIP) without requiring python-docx.

    DOCX is a ZIP archive with a Word/document.xml using OOXML namespace.
    mammoth reads this format directly, making python-docx optional for fixture
    generation.
    """
    import zipfile

    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml"'
        ' ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1"'
        ' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument"'
        ' Target="word/document.xml"/>'
        "</Relationships>"
    )
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:wpc="http://schemas.microsoft.com/office/word/2010/wordprocessingCanvas"'
        ' xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body>"
        f'<w:p><w:pPr><w:pStyle w:val="Heading1"/></w:pPr><w:r><w:t>{heading}</w:t></w:r></w:p>'
        f"<w:p><w:r><w:t>{body}</w:t></w:r></w:p>"
        "</w:body>"
        "</w:document>"
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("word/document.xml", document_xml)
    return buf.getvalue()


@pytest.fixture(scope="module")
def docx_bytes() -> bytes:
    """Generate a minimal real DOCX without requiring python-docx.

    Uses raw OOXML ZIP construction — mammoth reads this format directly.
    Falls back to python-docx if available (produces richer output).
    """
    try:
        import docx  # python-docx, only needed for richer fixture generation

        doc = docx.Document()
        doc.add_heading("Test Heading", level=1)
        doc.add_paragraph("Test paragraph content.")
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
    except ImportError:
        # Build a minimal valid DOCX from raw OOXML — no external deps needed
        return _build_minimal_docx("Test Heading", "Test paragraph content.")


@pytest.fixture(scope="module")
def xlsx_bytes() -> bytes:
    """Generate a minimal real XLSX with two sheets via openpyxl."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1.append(["Name", "Value"])
    ws1.append(["Alpha", 1])
    ws1.append(["Beta", 2])
    ws2 = wb.create_sheet("Sheet2")
    ws2.append(["X", "Y"])
    ws2.append(["10", "20"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture(scope="module")
def pptx_bytes() -> bytes:
    """Generate a minimal real PPTX with two slides via python-pptx."""
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # title + content layout
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "First Slide"
    slide1.placeholders[1].text = "Bullet one\nBullet two"
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Second Slide"
    slide2.placeholders[1].text = "Content here"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# OFFICE-01: DOCX extraction
# ---------------------------------------------------------------------------


class TestDocxExtraction:
    def test_returns_extraction_result(
        self, svc: OfficeExtractionService, docx_bytes: bytes
    ) -> None:
        result = svc.extract(docx_bytes, _DOCX_MIME, "test.docx")
        assert isinstance(result, ExtractionResult)

    def test_text_is_non_empty_string(
        self, svc: OfficeExtractionService, docx_bytes: bytes
    ) -> None:
        result = svc.extract(docx_bytes, _DOCX_MIME, "test.docx")
        assert isinstance(result.text, str)
        assert len(result.text) > 0

    def test_text_has_no_html_tags(self, svc: OfficeExtractionService, docx_bytes: bytes) -> None:
        """mammoth must be called in markdown mode, not HTML mode."""
        result = svc.extract(docx_bytes, _DOCX_MIME, "test.docx")
        assert "<p>" not in result.text
        assert "<h1>" not in result.text
        assert "<ul>" not in result.text

    def test_metadata_contains_word_count(
        self, svc: OfficeExtractionService, docx_bytes: bytes
    ) -> None:
        result = svc.extract(docx_bytes, _DOCX_MIME, "test.docx")
        assert "word_count" in result.metadata
        assert isinstance(result.metadata["word_count"], int)
        assert result.metadata["word_count"] > 0

    def test_result_is_frozen(self, svc: OfficeExtractionService, docx_bytes: bytes) -> None:
        result = svc.extract(docx_bytes, _DOCX_MIME, "test.docx")
        with pytest.raises((TypeError, AttributeError)):
            result.text = "mutated"  # type: ignore[misc]


# ---------------------------------------------------------------------------
# OFFICE-02: XLSX extraction
# ---------------------------------------------------------------------------


class TestXlsxExtraction:
    def test_returns_extraction_result(
        self, svc: OfficeExtractionService, xlsx_bytes: bytes
    ) -> None:
        result = svc.extract(xlsx_bytes, _XLSX_MIME, "test.xlsx")
        assert isinstance(result, ExtractionResult)

    def test_sheet_headers_present(self, svc: OfficeExtractionService, xlsx_bytes: bytes) -> None:
        result = svc.extract(xlsx_bytes, _XLSX_MIME, "test.xlsx")
        assert "## Sheet: Sheet1" in result.text
        assert "## Sheet: Sheet2" in result.text

    def test_gfm_table_format(self, svc: OfficeExtractionService, xlsx_bytes: bytes) -> None:
        """Each sheet produces GFM markdown table with | separator rows and --- divider."""
        result = svc.extract(xlsx_bytes, _XLSX_MIME, "test.xlsx")
        assert "| Name | Value |" in result.text
        assert "| --- | --- |" in result.text

    def test_metadata_sheet_count(self, svc: OfficeExtractionService, xlsx_bytes: bytes) -> None:
        result = svc.extract(xlsx_bytes, _XLSX_MIME, "test.xlsx")
        assert result.metadata["sheet_count"] == 2

    def test_row_cap_at_500(self, svc: OfficeExtractionService) -> None:
        """XLSX with 600 data rows is capped at 500 with truncation notice."""
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "BigSheet"
        ws.append(["Col"])
        for i in range(600):
            ws.append([str(i)])
        buf = io.BytesIO()
        wb.save(buf)
        big_xlsx = buf.getvalue()

        result = svc.extract(big_xlsx, _XLSX_MIME, "big.xlsx")
        # Count data rows in the markdown table (lines starting with |, excluding header and sep)
        table_rows = [line for line in result.text.splitlines() if line.startswith("|")]
        # header + separator + 500 data rows = 502
        assert len(table_rows) == 502
        assert "Truncated" in result.text


# ---------------------------------------------------------------------------
# OFFICE-03: PPTX extraction
# ---------------------------------------------------------------------------


class TestPptxExtraction:
    def test_returns_extraction_result(
        self, svc: OfficeExtractionService, pptx_bytes: bytes
    ) -> None:
        result = svc.extract(pptx_bytes, _PPTX_MIME, "test.pptx")
        assert isinstance(result, ExtractionResult)

    def test_slide_headers_present(self, svc: OfficeExtractionService, pptx_bytes: bytes) -> None:
        result = svc.extract(pptx_bytes, _PPTX_MIME, "test.pptx")
        assert "## Slide 1: First Slide" in result.text
        assert "## Slide 2: Second Slide" in result.text

    def test_bullet_content_present(self, svc: OfficeExtractionService, pptx_bytes: bytes) -> None:
        result = svc.extract(pptx_bytes, _PPTX_MIME, "test.pptx")
        assert "- Bullet one" in result.text or "Bullet one" in result.text

    def test_metadata_slide_count(self, svc: OfficeExtractionService, pptx_bytes: bytes) -> None:
        result = svc.extract(pptx_bytes, _PPTX_MIME, "test.pptx")
        assert result.metadata["slide_count"] == 2

    def test_untitled_slide_fallback(self, svc: OfficeExtractionService) -> None:
        """Slide with no title shape produces '(Untitled)' header."""
        from pptx import Presentation

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # blank layout
        prs.slides.add_slide(blank_layout)
        buf = io.BytesIO()
        prs.save(buf)
        result = svc.extract(buf.getvalue(), _PPTX_MIME, "untitled.pptx")
        assert "(Untitled)" in result.text


# ---------------------------------------------------------------------------
# OFFICE-04: MIME routing and error paths
# ---------------------------------------------------------------------------


class TestMimeRouting:
    def test_unsupported_mime_raises_value_error(
        self, svc: OfficeExtractionService, docx_bytes: bytes
    ) -> None:
        with pytest.raises(ValidationError, match="Unsupported"):
            svc.extract(docx_bytes, "application/octet-stream", "file.bin")

    def test_office_mime_types_contains_all_three(self) -> None:
        assert (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            in OfficeExtractionService.OFFICE_MIME_TYPES
        )
        assert (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            in OfficeExtractionService.OFFICE_MIME_TYPES
        )
        assert (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            in OfficeExtractionService.OFFICE_MIME_TYPES
        )

    def test_password_protected_returns_sentinel(self, svc: OfficeExtractionService) -> None:
        """Non-parseable bytes produce a sentinel result, not an exception."""
        # Garbage bytes that will fail parsing
        result = svc.extract(b"not a real office file at all", _DOCX_MIME, "broken.docx")
        assert isinstance(result, ExtractionResult)
        assert "password-protected" in result.text or "could not be extracted" in result.text
